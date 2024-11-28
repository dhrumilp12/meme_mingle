""" Helper functions for extracting text from various file types using Azure Form Recognizer """

"""Step 1: Import necessary modules"""
import os
import logging
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from io import BytesIO

# Import additional libraries
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

"""Step 2: Define the helper functions"""
# Define a function to get the Form Recognizer client
def get_form_recognizer_client():
    endpoint = os.getenv("AZURE_FORM_RECOGNIZER_ENDPOINT")
    key = os.getenv("AZURE_FORM_RECOGNIZER_KEY")
    return DocumentAnalysisClient(endpoint, AzureKeyCredential(key))


# Define a function to extract text from a file
def extract_text_from_file(file_content, file_mime_type):
    try:
        # Handle .docx files
        if file_mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            logging.info("Processing a .docx file")
            return extract_text_from_docx(file_content)
        
        # Handle .xlsx files
        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            logging.info("Processing an .xlsx file")
            return extract_text_from_xlsx(file_content)
        
        # Handle .pptx files
        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            logging.info("Processing a .pptx file")
            return extract_text_from_pptx(file_content)
        
        # Handle PDFs and images with Azure Form Recognizer
        else:
            client = get_form_recognizer_client()
            poller = client.begin_analyze_document(
                "prebuilt-read", document=file_content
            )
            result = poller.result()
            extracted_text = ""
            for page in result.pages:
                for line in page.lines:
                    extracted_text += line.content + "\n"
            return extracted_text

    except Exception as e:
        logging.error(f"Error extracting text from file: {e}")
        return ""


# Define a function to extract text from a .docx file
def extract_text_from_docx(file_content):
    extracted_text = ""
    try:
        document = DocxDocument(BytesIO(file_content))
        for para in document.paragraphs:
            extracted_text += para.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .docx file: {e}")
    return extracted_text


# Define a function to extract text from a .xlsx file
def extract_text_from_xlsx(file_content):
    extracted_text = ""
    try:
        workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                extracted_text += "\t".join(row_text) + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .xlsx file: {e}")
    return extracted_text


# Define a function to extract text from a .pptx file
def extract_text_from_pptx(file_content):
    extracted_text = ""
    try:
        presentation = Presentation(BytesIO(file_content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .pptx file: {e}")
    return extracted_text

def extract_text_from_txt(file_content):
    return file_content.decode('utf-8', errors='ignore')


# Define a function to check if a file type is supported
ALLOWED_MIME_TYPES = [
            'application/pdf',
            'image/jpeg',
            'image/png',
            'image/tiff',
            'image/bmp',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',        # .xlsx
            'application/vnd.openxmlformats-officedocument.presentationml.presentation', # .pptx
            'text/plain',  # .txt
            # Add more MIME types as needed
        ]