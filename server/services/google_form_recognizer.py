"""
Helper functions for extracting text from various file types using Google Cloud Document AI
"""

"""Step 1: Import necessary modules"""
import os
import logging
from google.cloud import documentai_v1 as documentai
from google.oauth2 import service_account
from io import BytesIO

# Import additional libraries
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

"""Step 2: Define the helper functions"""
# Define a function to get the Document AI client
def get_document_ai_client():
    """
    Initializes and returns a Document AI client.
    
    Returns:
        DocumentProcessorServiceClient: An authenticated Document AI client.
    """
    # Path to the service account key file
    service_account_path = os.getenv("GCP_DOCUMENT_AI_SERVICE_ACCOUNT")
    if not service_account_path:
        raise ValueError("GCP_DOCUMENT_AI_SERVICE_ACCOUNT environment variable not set.")
    
    # Credentials and client initialization
    credentials = service_account.Credentials.from_service_account_file(service_account_path)
    client = documentai.DocumentProcessorServiceClient(credentials=credentials)
    return client

# Define a function to extract text from a file using Google Document AI
def extract_text_from_file(file_content, file_mime_type):
    """
    Extracts text from a file based on its MIME type using Google Document AI or local libraries.

    Args:
        file_content (bytes): The content of the file.
        file_mime_type (str): The MIME type of the file.

    Returns:
        str: The extracted text.
    """
    try:
        # Handle .docx files
        if file_mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            logging.info("Processing a .docx file")
            return extract_text_from_docx(file_content)
        
        # Handle .xlsx files
        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            logging.info("Processing an .xlsx file")
            return extract_text_from_xlsx(file_content)
        
        # Handle .pptx files
        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            logging.info("Processing a .pptx file")
            return extract_text_from_pptx(file_content)
        
        # Handle PDFs and images with Google Document AI
        elif file_mime_type in ['application/pdf', 'image/jpeg', 'image/png', 'image/tiff', 'image/bmp']:
            logging.info("Processing a PDF or image file with Document AI")
            client = get_document_ai_client()
            project_id = os.getenv("GCP_PROJECT_ID")
            location = os.getenv("GCP_DOCUMENT_AI_LOCATION", "us")  # e.g., 'us'
            processor_id = os.getenv("GCP_DOCUMENT_AI_PROCESSOR_ID")  # The processor ID
            if not all([project_id, processor_id]):
                raise ValueError("GCP_PROJECT_ID and GCP_DOCUMENT_AI_PROCESSOR_ID must be set in environment variables.")
            
            # Construct the full resource name of the processor
            name = f"projects/{project_id}/locations/{location}/processors/{processor_id}"
            
            # Read the file content
            raw_document = documentai.RawDocument(content=file_content, mime_type=file_mime_type)
            
            # Configure the process request
            request = documentai.ProcessRequest(
                name=name,
                raw_document=raw_document
            )
            
            # Process the document
            result = client.process_document(request=request)
            
            # Extract text from the Document object
            extracted_text = result.document.text
            return extracted_text
        
        # Handle .txt files
        elif file_mime_type == 'text/plain':
            logging.info("Processing a .txt file")
            return extract_text_from_txt(file_content)
        
        else:
            logging.warning(f"Unsupported MIME type: {file_mime_type}. Returning empty string.")
            return ""
    
    except Exception as e:
        logging.error(f"Error extracting text from file: {e}")
        return ""

# Define a function to extract text from a .docx file
def extract_text_from_docx(file_content):
    """
    Extracts text from a .docx file.

    Args:
        file_content (bytes): The content of the .docx file.

    Returns:
        str: The extracted text.
    """
    extracted_text = ""
    try:
        document = DocxDocument(BytesIO(file_content))
        for para in document.paragraphs:
            extracted_text += para.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .docx file: {e}")
    return extracted_text

# Define a function to extract text from a .xlsx file
def extract_text_from_xlsx(file_content):
    """
    Extracts text from a .xlsx file.

    Args:
        file_content (bytes): The content of the .xlsx file.

    Returns:
        str: The extracted text.
    """
    extracted_text = ""
    try:
        workbook = openpyxl.load_workbook(BytesIO(file_content), data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                extracted_text += "\t".join(row_text) + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .xlsx file: {e}")
    return extracted_text

# Define a function to extract text from a .pptx file
def extract_text_from_pptx(file_content):
    """
    Extracts text from a .pptx file.

    Args:
        file_content (bytes): The content of the .pptx file.

    Returns:
        str: The extracted text.
    """
    extracted_text = ""
    try:
        presentation = Presentation(BytesIO(file_content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .pptx file: {e}")
    return extracted_text

# Define a function to extract text from a .txt file
def extract_text_from_txt(file_content):
    """
    Extracts text from a .txt file.

    Args:
        file_content (bytes): The content of the .txt file.

    Returns:
        str: The extracted text.
    """
    try:
        return file_content.decode('utf-8', errors='ignore')
    except Exception as e:
        logging.error(f"Error extracting text from .txt file: {e}")
        return ""

# Define a function to check if a file type is supported
ALLOWED_MIME_TYPES = [
    'application/pdf',
    'image/jpeg',
    'image/png',
    'image/tiff',
    'image/bmp',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',        # .xlsx
    'application/vnd.openxmlformats-officedocument.presentationml.presentation', # .pptx
    'text/plain',  # .txt
    'application/msword',  # .doc
    'application/vnd.ms-excel',  # .xls
    'application/vnd.ms-powerpoint',  # .ppt
    'application/zip',  # .zip
    # Add more MIME types as needed
    
]
